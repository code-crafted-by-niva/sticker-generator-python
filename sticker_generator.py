# pip install python-pptx
# pip install pptxtopdf

from pptx import Presentation
# from pptxtopdf import convert
import pandas as pd
import os

awardFolder = 'output\\'
template = 'Trophy stickers.pptx'
excelFile = 'Trophy sticker -onsite (1).xlsx'
search_str_1 = '{{Name}}'
search_str_2 = '{{Category}}'

df = pd.read_excel(excelFile)
new_ppt = Presentation()
for index, row in df.iterrows():
    ppt = Presentation(template)
    
    print('\n\nExecuting for :', row['Awardee Name'])
    replace_str_1 = row['Awardee Name']
    replace_str_2 = row['Award Category']
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # run.text = run.text.replace(search_str_1, replace_str_1).replace(search_str_2, replace_str_2)
                        print(run.text)
                        if(run.text.find(search_str_1))!=-1:
                            run.text = run.text.replace(search_str_1, replace_str_1)
                        if(run.text.find(search_str_2))!=-1:
                            run.text = run.text.replace(search_str_2, replace_str_2)
        # new_slide = new_ppt.slides.add_slide(slide.slide_layout)
    tempPPT = awardFolder+row['Awardee Name']+'.pptx'
    ppt.save(tempPPT)

# new_ppt.save("output.pptx")