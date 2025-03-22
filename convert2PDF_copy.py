# pip install python-pptx
# pip install pptxtopdf

from pptx import Presentation
from pptxtopdf import convert
import pandas as pd
import os

awardFolder = 'onsite' + "\\"
template = 'TROPHY STICKERS  - Copy.pptx'
excelFile = 'Trophy sticker -onsite.xlsx'
search_str_1 = '{{Name}}'
search_str_2 = '{{Category}}'

df = pd.read_excel(excelFile, sheet_name='Onsite')
new_ppt = Presentation()
for index, row in df.iterrows():
    ppt = Presentation(template)
    
    print('\n\nExecuting for :', row['Nominee Name'])
    replace_str_1 = row['Nominee Name']
    replace_str_2 = row['Award Category']
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # run.text = run.text.replace(search_str_1, replace_str_1).replace(search_str_2, replace_str_2)
                        print(run.text)
                        if(run.text.find(search_str_1))!=-1:
                            run.text = run.text.replace(search_str_1, replace_str_1)
                        if(run.text.find(search_str_2))!=-1:
                            run.text = run.text.replace(search_str_2, replace_str_2)


    tempPPT = awardFolder+row['Nominee Name']+'.pptx'
    ppt.save(tempPPT)
    convert(tempPPT, awardFolder)
    os.remove(tempPPT)

    

